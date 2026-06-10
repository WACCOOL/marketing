#!/usr/bin/env python3
"""PPT deck builder (PRD §8).

Reads a build spec as JSON on STDIN (written by apps/generator/src/ppt.ts):

  {"templatePath": str, "layoutMap": {canonical: layoutName}, "outPath": str,
   "slides": [{"id": str, "layout": canonical,
               "fields": {"title"?, "subtitle"?, "bullets"?: [str], "body"?,
                          "body2"?,
                          "images"?: [{"path", "width", "height", "caption"?}],
                          "table"?: {"headers": [str], "rows": [[str]]},
                          "quote"?: {"text": str, "attribution"?: str},
                          "chart"?: {"chartType": "column"|"bar"|"line"|"pie",
                                     "categories": [str],
                                     "series": [{"name": str, "values": [num]}]},
                          "items"?: [str],            # diagram / process
                          "video"?: {"path": str, "mimeType": str,
                                     "caption"?: str}}}]}

Canonical-layout specifics:
  - agenda:  bullets get arabic auto-numbering ("1." "2." …) via raw pPr XML.
  - quote:   quote.text -> first body placeholder; attribution ("— …") -> the
             next body placeholder, else a textbox below the quote area.
  - chart:   a NATIVE pptx chart over the body/picture/table placeholder area;
             colors/fonts inherit from the template theme.
  - diagram: items as ROUNDED_RECTANGLE shapes in a centered grid (<=3 per
             row) filled with the theme's ACCENT_1.
  - process: items as a left-to-right PENTAGON+CHEVRON sequence (two rows
             when more than 6), same ACCENT_1 theme fill.
  - video:   embedded movie sized 16:9 over the media/picture placeholder.
  - title_content_image: body+bullets share the body placeholder; the first
             image fills the picture placeholder (placeholder cover-crop),
             else a right-half positioned picture.

Unknown field keys (notably `imagePrompt`, and `prompt` on image entries) are
ignored — only the keys named above are read.

Opens the template, drops all its existing slides (masters/layouts/theme stay —
that's what keeps exports on-brand), then adds one slide per spec entry and
fills placeholders by classified type. Body-type content (body, then bullets,
then body2) each consume the next free body placeholder; when none is left the
content is appended as extra paragraphs to the LAST body placeholder used (so
"Title and Content" merges body + bullets into its single body), and only a
slide with no body placeholder at all falls back to a positioned textbox.
When a mapped layout lacks any other needed placeholder we never error: a
positioned textbox/picture/shape is added at sensible margins and a warning is
appended to warnings[].

IMPORTANT: this script NEVER sets fonts, colors, or sizes on text — everything
inherits from the template. (Diagram/process shapes get a THEME fill —
ACCENT_1 — which is the template's own brand color, not a hardcoded one.)

Output: {"ok": true, "pptxPath": outPath, "warnings": [str]}
        or {"ok": false, "error": {"code": str, "message": str}} with exit 1.
"""

import json
import sys

from pptlib import build_layout_index, classify_placeholder, resolve_layout

# Fallback geometry (fractions of the slide) used when the mapped layout lacks
# a needed placeholder. Title across the top ~10%, content stacked below.
FALLBACK = {
    "title": (0.05, 0.03, 0.90, 0.10),
    "subtitle": (0.05, 0.14, 0.90, 0.08),
    "body": (0.05, 0.25, 0.90, 0.32),
    "body_lower": (0.05, 0.58, 0.90, 0.32),
    "picture": (0.10, 0.20, 0.80, 0.62),
    # title_content_image: text keeps the left half, image takes the right.
    "picture_right": (0.52, 0.20, 0.43, 0.62),
    "table": (0.05, 0.25, 0.90, 0.60),
    "caption": (0.05, 0.90, 0.90, 0.08),
    # quote: a large box at mid-slide, attribution just below it.
    "quote": (0.10, 0.32, 0.80, 0.28),
    "attribution": (0.10, 0.62, 0.80, 0.08),
    # General content area below the title for chart/diagram/process/video.
    "content": (0.05, 0.20, 0.90, 0.70),
}

# chartType -> XL_CHART_TYPE member name (resolved lazily inside add_slide_chart
# so the script keeps its import-on-use pattern).
CHART_TYPE_NAMES = {
    "column": "COLUMN_CLUSTERED",
    "bar": "BAR_CLUSTERED",
    "line": "LINE",
    "pie": "PIE",
}


def fail(code, message):
    print(json.dumps({"ok": False, "error": {"code": code, "message": message}}))
    sys.exit(1)


def remove_all_slides(prs):
    """Drop every existing slide (sldIdLst entries + their rels)."""
    from pptx.oxml.ns import qn

    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001 — no public API for this
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.get(qn("r:id")))
        sld_id_lst.remove(sld_id)


def bucket_placeholders(slide):
    """Slide placeholders grouped by classified type, in idx order."""
    buckets = {
        "title": [],
        "subtitle": [],
        "body": [],
        "picture": [],
        "table": [],
        "media": [],
    }
    for ph in slide.placeholders:
        kind = classify_placeholder(ph.placeholder_format)
        if kind in buckets:
            buckets[kind].append(ph)
    for phs in buckets.values():
        phs.sort(key=lambda p: p.placeholder_format.idx)
    return buckets


def write_frame(tf, text=None, bullets=None, fresh=True):
    """Write content into a text frame; returns the paragraphs it created.

    fresh=True replaces the frame's content (placeholder prompt text included);
    fresh=False appends paragraphs after whatever is already there — that's the
    body+bullets merge path. Formatting is always inherited.
    """
    if bullets is not None:
        lines = [str(b) for b in bullets]
    else:
        lines = str(text).split("\n")
    created = []
    if fresh:
        tf.text = lines[0]
        para = tf.paragraphs[0]
        if bullets is not None:
            para.level = 0
        created.append(para)
        lines = lines[1:]
    for line in lines:
        para = tf.add_paragraph()
        para.text = line
        if bullets is not None:
            para.level = 0
        created.append(para)
    return created


def set_text(text_frame, text):
    """Fill a text frame, one paragraph per line. Formatting is inherited."""
    write_frame(text_frame, text=text)


def add_textbox(slide, sw, sh, box, text, bullets=None):
    """Positioned-textbox fallback for a missing placeholder."""
    from pptx.util import Emu

    x, y, w, h = box
    shape = slide.shapes.add_textbox(
        Emu(int(sw * x)), Emu(int(sh * y)), Emu(int(sw * w)), Emu(int(sh * h))
    )
    tf = shape.text_frame
    tf.word_wrap = True
    return write_frame(tf, text=text, bullets=bullets)


def take_area(buckets, kinds):
    """Pop the first placeholder of the given kinds that has resolvable
    geometry and return its (left, top, width, height) in EMU, else None.

    Placeholders without explicit/inherited geometry stay in their bucket so
    later content can still fill them with text.
    """
    for kind in kinds:
        phs = buckets.get(kind) or []
        for i, ph in enumerate(phs):
            if (
                ph.left is not None
                and ph.top is not None
                and ph.width is not None
                and ph.height is not None
            ):
                phs.pop(i)
                return (int(ph.left), int(ph.top), int(ph.width), int(ph.height))
    return None


def fallback_area(sw, sh, key):
    x, y, w, h = FALLBACK[key]
    return (int(sw * x), int(sh * y), int(sw * w), int(sh * h))


def add_picture_fit(slide, image, sw, sh, box):
    """Positioned-picture fallback: fit inside the box, centered, no crop."""
    from pptx.util import Emu

    x, y, w, h = box
    box_w, box_h = sw * w, sh * h
    ar = image["width"] / max(image["height"], 1)
    if box_w / box_h > ar:
        draw_h, draw_w = box_h, box_h * ar
    else:
        draw_w, draw_h = box_w, box_w / ar
    left = sw * x + (box_w - draw_w) / 2
    top = sh * y + (box_h - draw_h) / 2
    return slide.shapes.add_picture(
        image["path"], Emu(int(left)), Emu(int(top)), Emu(int(draw_w)), Emu(int(draw_h))
    )


def add_picture_cover(slide, image, sw, sh):
    """image_full: scale to cover the whole slide, center, crop the overflow.

    The picture is stretched to the slide frame and the source overflow is
    removed via crop fractions, so the visible region keeps the image's aspect.
    """
    from pptx.util import Emu

    pic = slide.shapes.add_picture(image["path"], 0, 0, Emu(int(sw)), Emu(int(sh)))
    ar = image["width"] / max(image["height"], 1)
    sar = sw / sh
    if ar > sar:
        crop = (1 - sar / ar) / 2  # image is wider: trim left/right
        pic.crop_left = crop
        pic.crop_right = crop
    elif ar < sar:
        crop = (1 - ar / sar) / 2  # image is taller: trim top/bottom
        pic.crop_top = crop
        pic.crop_bottom = crop
    return pic


def apply_auto_numbering(paragraphs, warn):
    """Arabic "1." auto-numbering on agenda paragraphs.

    python-pptx has no bullet/numbering API, so this sets <a:buAutoNum
    type="arabicPeriod"/> on each paragraph's pPr via lxml, replacing any
    inherited buChar/buNone. The element is inserted before tabLst/defRPr/
    extLst to keep the a:pPr child order schema-valid. If the XML manipulation
    fails we fall back to literal "1. " text prefixes plus a warning.
    """
    try:
        from pptx.oxml.ns import qn

        successors = (qn("a:tabLst"), qn("a:defRPr"), qn("a:extLst"))
        for para in paragraphs:
            p_pr = para._p.get_or_add_pPr()  # noqa: SLF001 — no public API
            for tag in ("a:buNone", "a:buAutoNum", "a:buChar"):
                for el in p_pr.findall(qn(tag)):
                    p_pr.remove(el)
            bu = p_pr.makeelement(qn("a:buAutoNum"), {"type": "arabicPeriod"})
            anchor = next((c for c in p_pr if c.tag in successors), None)
            if anchor is not None:
                anchor.addprevious(bu)
            else:
                p_pr.append(bu)
    except Exception as e:  # noqa: BLE001 — numbering is cosmetic, never fatal
        warn(f"agenda auto-numbering failed ({e}); used literal number prefixes")
        for i, para in enumerate(paragraphs, start=1):
            para.text = f"{i}. {para.text}"


def add_slide_quote(slide, quote, buckets, sw, sh, warn):
    """quote.text into the first body placeholder (else a large mid-slide
    textbox); attribution as "— …" into the next body placeholder (else a
    textbox below the quote area — by design, not a warning)."""
    if buckets["body"]:
        write_frame(buckets["body"].pop(0).text_frame, text=quote.get("text", ""))
    else:
        warn("layout has no body placeholder for the quote; added a positioned textbox")
        add_textbox(slide, sw, sh, FALLBACK["quote"], quote.get("text", ""))

    attribution = quote.get("attribution")
    if attribution:
        text = f"— {attribution}"
        if buckets["body"]:
            write_frame(buckets["body"].pop(0).text_frame, text=text)
        else:
            add_textbox(slide, sw, sh, FALLBACK["attribution"], text)


def add_slide_chart(slide, spec_chart, buckets, sw, sh, warn):
    """A NATIVE chart over the body/picture/table placeholder area (else
    margins below the title). No styling — charts inherit the template theme."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Emu

    chart_type = spec_chart.get("chartType", "column")
    member = CHART_TYPE_NAMES.get(chart_type)
    if member is None:
        warn(f"unknown chartType '{chart_type}'; rendered a clustered column chart")
        member = CHART_TYPE_NAMES["column"]
    xl_type = getattr(XL_CHART_TYPE, member)

    series = list(spec_chart.get("series") or [])
    if chart_type == "pie" and len(series) > 1:
        warn(f"pie charts plot a single series; dropped {len(series) - 1} extra series")
        series = series[:1]

    data = CategoryChartData()
    data.categories = [str(c) for c in spec_chart.get("categories") or []]
    for s in series:
        data.add_series(str(s.get("name", "")), tuple(float(v) for v in s.get("values") or []))

    area = take_area(buckets, ("body", "picture", "table"))
    if area is None:
        warn("no body/picture/table placeholder; placed chart at default margins")
        area = fallback_area(sw, sh, "content")
    x, y, w, h = area
    slide.shapes.add_chart(xl_type, Emu(x), Emu(y), Emu(w), Emu(h), data)


def add_theme_shape(slide, shape_type, x, y, w, h, text):
    """An autoshape filled with the theme's ACCENT_1 so brand colors apply.
    Text styling stays inherited; word wrap on so labels stay inside."""
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.util import Emu

    shape = slide.shapes.add_shape(shape_type, Emu(x), Emu(y), Emu(w), Emu(h))
    shape.fill.solid()
    shape.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    tf = shape.text_frame
    tf.word_wrap = True
    write_frame(tf, text=text)
    return shape


def add_slide_diagram(slide, items, buckets, sw, sh, warn):
    """items as ROUNDED_RECTANGLE shapes in a centered grid (<=3 per row),
    evenly spaced within the body placeholder area (else default margins)."""
    from pptx.enum.shapes import MSO_SHAPE

    area = take_area(buckets, ("body",))
    if area is None:
        warn("no body placeholder; placed diagram at default margins")
        area = fallback_area(sw, sh, "content")
    x, y, w, h = area

    n = len(items)
    cols = min(3, n)
    rows = (n + cols - 1) // cols
    gap_x = int(w * 0.05) if cols > 1 else 0
    gap_y = int(h * 0.08) if rows > 1 else 0
    cell_w = (w - gap_x * (cols - 1)) // cols
    cell_h = (h - gap_y * (rows - 1)) // rows

    for i, item in enumerate(items):
        r, c = divmod(i, cols)
        in_row = cols if (r < rows - 1 or n % cols == 0) else n % cols
        row_w = in_row * cell_w + (in_row - 1) * gap_x
        left = x + (w - row_w) // 2 + c * (cell_w + gap_x)
        top = y + r * (cell_h + gap_y)
        add_theme_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, cell_w, cell_h, str(item))


def add_slide_process(slide, items, buckets, sw, sh, warn):
    """items as a left-to-right chevron sequence (PENTAGON for step 1, CHEVRON
    after), wrapping onto a second row past 6 items. Theme ACCENT_1 fill."""
    from pptx.enum.shapes import MSO_SHAPE

    area = take_area(buckets, ("body",))
    if area is None:
        warn("no body placeholder; placed process steps at default margins")
        area = fallback_area(sw, sh, "content")
    x, y, w, h = area

    n = len(items)
    rows = 2 if n > 6 else 1
    per_row = (n + rows - 1) // rows
    gap_x = int(w * 0.02) if per_row > 1 else 0
    gap_y = int(h * 0.10) if rows > 1 else 0
    cell_w = (w - gap_x * (per_row - 1)) // per_row
    # Chevrons read as a strip, not a panel: cap their height and center the
    # row block vertically inside the area.
    cell_h = min((h - gap_y * (rows - 1)) // rows, int(sh * 0.20))
    block_h = rows * cell_h + (rows - 1) * gap_y
    top0 = y + max((h - block_h) // 2, 0)

    for i, item in enumerate(items):
        r, c = divmod(i, per_row)
        shape_type = MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON
        left = x + c * (cell_w + gap_x)
        top = top0 + r * (cell_h + gap_y)
        add_theme_shape(slide, shape_type, left, top, cell_w, cell_h, str(item))


def add_slide_video(slide, video, buckets, sw, sh, warn):
    """Embed the (already-localized) movie file sized 16:9 over the media (or
    picture) placeholder area, else centered below the title. The largest 16:9
    box that fits the area is used, centered; the poster frame stays the
    python-pptx default. Caption into a body placeholder or a textbox below."""
    from pptx.util import Emu

    area = take_area(buckets, ("media", "picture"))
    if area is None:
        warn("no media/picture placeholder; placed video at default margins")
        area = fallback_area(sw, sh, "content")
    x, y, w, h = area

    if w * 9 > h * 16:  # area wider than 16:9 — height limits
        vid_h, vid_w = h, int(h * 16 / 9)
    else:
        vid_w, vid_h = w, int(w * 9 / 16)
    left = x + (w - vid_w) // 2
    top = y + (h - vid_h) // 2
    slide.shapes.add_movie(
        video["path"],
        Emu(left),
        Emu(top),
        Emu(vid_w),
        Emu(vid_h),
        mime_type=video.get("mimeType", "video/mp4"),
    )

    caption = video.get("caption")
    if caption:
        if buckets["body"]:
            write_frame(buckets["body"].pop(0).text_frame, text=caption)
        else:
            # By design: drop the caption just below the embedded movie.
            box = (
                left / sw,
                min((top + vid_h) / sh, 0.92),
                vid_w / sw,
                0.07,
            )
            add_textbox(slide, sw, sh, box, caption)


def fill_table_cells(table, headers, rows):
    """Header row + data rows. Cell text only — styling stays the template's."""
    for c, header in enumerate(headers):
        table.cell(0, c).text = str(header)
    for r, row in enumerate(rows, start=1):
        for c in range(len(headers)):
            table.cell(r, c).text = str(row[c]) if c < len(row) else ""


def add_slide_table(slide, spec_table, buckets, sw, sh, warn):
    """Table via the table placeholder, else a body placeholder's area, else
    fallback margins (with a warning)."""
    from pptx.util import Emu

    headers = spec_table["headers"]
    rows = spec_table.get("rows", [])
    n_rows, n_cols = len(rows) + 1, len(headers)

    if buckets["table"]:
        ph = buckets["table"].pop(0)
        graphic_frame = ph.insert_table(rows=n_rows, cols=n_cols)
        fill_table_cells(graphic_frame.table, headers, rows)
        return

    geom = take_area(buckets, ("body",))
    if geom is None:
        warn("no table/body placeholder; placed table at default margins")
        geom = fallback_area(sw, sh, "table")
    x, y, w, h = geom
    graphic_frame = slide.shapes.add_table(n_rows, n_cols, Emu(x), Emu(y), Emu(w), Emu(h))
    fill_table_cells(graphic_frame.table, headers, rows)


def fill_slide(slide, spec_slide, sw, sh, warn):
    canonical = spec_slide["layout"]
    fields = spec_slide.get("fields", {})
    buckets = bucket_placeholders(slide)

    # --- title / subtitle ----------------------------------------------------
    if fields.get("title"):
        if buckets["title"]:
            set_text(buckets["title"].pop(0).text_frame, fields["title"])
        else:
            warn("layout has no title placeholder; added a positioned textbox")
            add_textbox(slide, sw, sh, FALLBACK["title"], fields["title"])

    if fields.get("subtitle"):
        if buckets["subtitle"]:
            set_text(buckets["subtitle"].pop(0).text_frame, fields["subtitle"])
        elif buckets["body"]:
            # Common in non-title layouts: the subtitle area is a body placeholder.
            set_text(buckets["body"].pop(0).text_frame, fields["subtitle"])
        else:
            warn("layout has no subtitle placeholder; added a positioned textbox")
            add_textbox(slide, sw, sh, FALLBACK["subtitle"], fields["subtitle"])

    # --- quote ----------------------------------------------------------------
    if fields.get("quote") and fields["quote"].get("text"):
        add_slide_quote(slide, fields["quote"], buckets, sw, sh, warn)

    # --- body-type content (body, then bullets, then body2) --------------------
    # Each consumes the next free body placeholder; once none is left, the
    # content merges into the LAST body placeholder used as extra paragraphs
    # (e.g. "Title and Content" gets body paragraphs first, bullets after, in
    # its single body). Positioned textboxes only when the layout has no body
    # placeholder at all.
    last_body_tf = None
    body_fallback_used = 0

    def fill_body(text=None, bullets=None, label="body"):
        nonlocal last_body_tf, body_fallback_used
        if buckets["body"]:
            tf = buckets["body"].pop(0).text_frame
            last_body_tf = tf
            return write_frame(tf, text=text, bullets=bullets)
        if last_body_tf is not None:
            return write_frame(last_body_tf, text=text, bullets=bullets, fresh=False)
        warn(f"layout has no body placeholder for {label}; added a positioned textbox")
        box = FALLBACK["body_lower"] if body_fallback_used else FALLBACK["body"]
        body_fallback_used += 1
        return add_textbox(slide, sw, sh, box, text, bullets=bullets)

    if fields.get("body"):
        fill_body(text=fields["body"], label="body")
    if fields.get("bullets"):
        paras = fill_body(bullets=fields["bullets"], label="bullets")
        if canonical == "agenda":
            apply_auto_numbering(paras, warn)
    if fields.get("body2"):
        fill_body(text=fields["body2"], label="body2")

    # --- chart / diagram / process ---------------------------------------------
    if fields.get("chart"):
        add_slide_chart(slide, fields["chart"], buckets, sw, sh, warn)

    if fields.get("items"):
        if canonical == "process":
            add_slide_process(slide, fields["items"], buckets, sw, sh, warn)
        else:
            add_slide_diagram(slide, fields["items"], buckets, sw, sh, warn)

    # --- images ---------------------------------------------------------------
    images = list(fields.get("images") or [])
    if images and canonical == "image_full":
        # Full-bleed cover crop for the first image; extras go through the
        # generic placeholder path below.
        image = images.pop(0)
        add_picture_cover(slide, image, sw, sh)
        if image.get("caption"):
            if buckets["body"]:
                set_text(buckets["body"].pop(0).text_frame, image["caption"])
            else:
                # By design (not a fallback warning): full-bleed layouts rarely
                # carry a caption placeholder.
                add_textbox(slide, sw, sh, FALLBACK["caption"], image["caption"])

    for image in images:
        ph = buckets["picture"].pop(0) if buckets["picture"] else None
        if ph is not None and hasattr(ph, "insert_picture"):
            # The picture placeholder cover-crops to its own frame.
            ph.insert_picture(image["path"])
        elif canonical == "title_content_image":
            warn("layout has no free picture placeholder; placed image on the right half")
            add_picture_fit(slide, image, sw, sh, FALLBACK["picture_right"])
        else:
            warn("layout has no free picture placeholder; placed image at default margins")
            add_picture_fit(slide, image, sw, sh, FALLBACK["picture"])
        if image.get("caption"):
            if buckets["body"]:
                set_text(buckets["body"].pop(0).text_frame, image["caption"])
            else:
                warn("layout has no free placeholder for the image caption; added a positioned textbox")
                add_textbox(slide, sw, sh, FALLBACK["caption"], image["caption"])

    # --- video ------------------------------------------------------------------
    if fields.get("video") and fields["video"].get("path"):
        add_slide_video(slide, fields["video"], buckets, sw, sh, warn)

    # --- table ----------------------------------------------------------------
    if fields.get("table") and fields["table"].get("headers"):
        add_slide_table(slide, fields["table"], buckets, sw, sh, warn)


def main():
    try:
        spec = json.load(sys.stdin)
    except Exception as e:  # noqa: BLE001
        fail("bad_spec", f"could not parse build spec JSON: {e}")
    for key in ("templatePath", "layoutMap", "outPath", "slides"):
        if key not in spec:
            fail("bad_spec", f"build spec is missing '{key}'")

    try:
        from pptx import Presentation

        prs = Presentation(spec["templatePath"])
    except Exception as e:  # noqa: BLE001
        fail("open_failed", f"could not open template: {e}")

    layout_index = build_layout_index(prs)
    layout_map = spec["layoutMap"]
    warnings = []

    try:
        remove_all_slides(prs)
        sw, sh = int(prs.slide_width), int(prs.slide_height)

        for spec_slide in spec["slides"]:
            canonical = spec_slide["layout"]
            layout_name = layout_map.get(canonical)
            if not layout_name:
                fail("layout_unmapped", f"no template mapping for layout '{canonical}'")
            layout = resolve_layout(layout_index, layout_name)
            if layout is None:
                fail(
                    "layout_not_found",
                    f"template has no slide layout named '{layout_name}' (mapped from '{canonical}')",
                )

            slide = prs.slides.add_slide(layout)
            fill_slide(
                slide,
                spec_slide,
                sw,
                sh,
                lambda msg, sid=spec_slide["id"]: warnings.append(f"slide {sid}: {msg}"),
            )

        prs.save(spec["outPath"])
    except SystemExit:
        raise
    except Exception as e:  # noqa: BLE001
        fail("build_failed", f"deck build failed: {e}")

    print(json.dumps({"ok": True, "pptxPath": spec["outPath"], "warnings": warnings}))


if __name__ == "__main__":
    main()
