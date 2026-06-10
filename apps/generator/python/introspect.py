#!/usr/bin/env python3
"""PPT template introspection (PRD §8).

Usage: python3 introspect.py <template.pptx>

Prints ONE JSON object to stdout describing every slide layout (across all
masters) and a suggested canonical->layout-name map. The shape is a contract
shared with apps/generator/src/ppt.ts and the PPT Templates admin UI:

  {"ok": true, "slideWidthEmu": int, "slideHeightEmu": int,
   "layouts": [{"index": i, "name": str,
                "placeholders": [{"idx": int, "type": str, "name": str,
                                  "xEmu": int, "yEmu": int, "wEmu": int, "hEmu": int}]}],
   "suggestedMap": {"<canonical>": "<layout name>", ...}}

Errors: {"ok": false, "error": {"code": str, "message": str}} and exit 1.
"""

import json
import sys

from pptlib import build_layout_index, classify_placeholder

# Canonical deck layouts — keep in lockstep with packages/shared/src/ppt.ts.
CANONICAL_LAYOUTS = [
    "title",
    "title_content",
    "title_content_image",
    "two_column",
    "image_full",
    "image_caption",
    "agenda",
    "quote",
    "chart",
    "diagram",
    "process",
    "video",
    "table",
    "section",
]

# A canonical only lands in suggestedMap when its best layout scores at least
# this much — unsure matches are omitted so admins map them by hand.
MIN_SCORE = 3


def fail(code, message):
    print(json.dumps({"ok": False, "error": {"code": code, "message": message}}))
    sys.exit(1)


def emu(value):
    """Placeholder geometry can be inherited (None); report 0 in that case."""
    return int(value) if value is not None else 0


def placeholder_counts(layout):
    counts = {
        "title": 0,
        "subtitle": 0,
        "body": 0,
        "picture": 0,
        "table": 0,
        "media": 0,
        "other": 0,
    }
    for ph in layout.placeholders:
        counts[classify_placeholder(ph.placeholder_format)] += 1
    return counts


def score_layout(canonical, name, counts):
    """Heuristic confidence that a template layout serves a canonical role.

    Two signals: the layout's own name (PowerPoint templates keep conventional
    names like "Title Slide" / "Section Header" / "Two Content") and its
    placeholder composition. Either alone can clear MIN_SCORE.
    """
    n = name.lower()
    has_text = counts["title"] or counts["subtitle"] or counts["body"]
    empty = not has_text and not counts["picture"] and not counts["table"]
    score = 0

    if canonical == "title":
        if "title slide" in n or n == "title":
            score += 4
        if counts["title"] and counts["subtitle"] and not counts["body"]:
            score += 3
    elif canonical == "title_content":
        if "title and content" in n:
            score += 4
        elif "content" in n and "two" not in n:
            score += 2
        if counts["title"] and counts["body"] == 1 and not counts["picture"] and not counts["table"]:
            score += 3
    elif canonical == "title_content_image":
        if (
            "content and image" in n
            or "text and image" in n
            or "content with image" in n
        ):
            score += 4
        elif ("picture" in n or "image" in n) and ("content" in n or "text" in n):
            score += 2
        if counts["title"] and counts["body"] and counts["picture"]:
            score += 3
    elif canonical == "two_column":
        if "two" in n or "comparison" in n:
            score += 4
        if counts["title"] and counts["body"] == 2:
            score += 3
    elif canonical == "image_full":
        if "blank" in n and empty:
            score += 3
        if ("picture" in n or "image" in n or "photo" in n) and "caption" not in n:
            score += 2
        if counts["picture"] and not has_text:
            score += 3
        elif empty:
            score += 1  # cover-crop works on a bare canvas too
    elif canonical == "image_caption":
        if "caption" in n:
            score += 4
        elif "picture" in n or "image" in n or "photo" in n:
            score += 1
        if counts["picture"] and has_text:
            score += 3
    elif canonical == "agenda":
        if "agenda" in n:
            score += 4
        if counts["title"] and counts["body"]:
            score += 1  # weak: any title+body layout can host a numbered list
    elif canonical == "quote":
        if "quote" in n or "quotation" in n:
            score += 4
        if counts["body"] and not counts["picture"] and not counts["table"]:
            score += 1  # weak: a body area can host the pull quote
    elif canonical == "chart":
        if "chart" in n or "graph" in n:
            score += 4
        if counts["title"] and counts["body"] and not counts["picture"]:
            score += 1  # weak: the chart is drawn over a body area
    elif canonical == "diagram":
        if "diagram" in n:
            score += 4
        if counts["title"] and counts["body"]:
            score += 1  # weak: shapes are drawn over a body area
    elif canonical == "process":
        if "process" in n or "flow" in n or "steps" in n or "timeline" in n:
            score += 4
        if counts["title"] and counts["body"]:
            score += 1  # weak: chevrons are drawn over a body area
    elif canonical == "video":
        if "video" in n or "media" in n:
            score += 4
        if counts["media"]:
            score += 4  # a media placeholder is a definitive signal
        elif counts["picture"]:
            score += 1  # weak: the movie can sit over a picture area
    elif canonical == "table":
        if "table" in n:
            score += 4
        if counts["table"]:
            score += 4
        elif counts["title"] and counts["body"]:
            score += 1  # weak fallback: a body area can host an added table
    elif canonical == "section":
        if "section" in n:
            score += 4
        if (
            counts["title"]
            and not counts["body"]
            and not counts["subtitle"]
            and not counts["picture"]
            and not counts["table"]
        ):
            score += 2

    return score


def suggest_map(layouts):
    """Best-scoring layout per canonical; omit anything below MIN_SCORE."""
    suggested = {}
    for canonical in CANONICAL_LAYOUTS:
        best_name = None
        best_score = 0
        for entry in layouts:
            s = score_layout(canonical, entry["name"], entry["counts"])
            if s > best_score:
                best_score = s
                best_name = entry["name"]
        if best_name is not None and best_score >= MIN_SCORE:
            suggested[canonical] = best_name
    return suggested


def main():
    if len(sys.argv) != 2:
        fail("usage", "usage: introspect.py <template.pptx>")

    try:
        from pptx import Presentation

        prs = Presentation(sys.argv[1])
    except Exception as e:  # noqa: BLE001 — surface any open failure as JSON
        fail("open_failed", f"could not open template: {e}")

    try:
        # Enumerate every layout across every master; the name index (first
        # name wins) decides which one a duplicate name resolves to, but the
        # admin UI should still see them all.
        layouts = []
        index = 0
        for master in prs.slide_masters:
            for layout in master.slide_layouts:
                placeholders = []
                for ph in layout.placeholders:
                    fmt = ph.placeholder_format
                    placeholders.append(
                        {
                            "idx": int(fmt.idx),
                            "type": classify_placeholder(fmt),
                            "name": ph.name or "",
                            "xEmu": emu(ph.left),
                            "yEmu": emu(ph.top),
                            "wEmu": emu(ph.width),
                            "hEmu": emu(ph.height),
                        }
                    )
                layouts.append(
                    {
                        "index": index,
                        "name": layout.name or "",
                        "placeholders": placeholders,
                        "counts": placeholder_counts(layout),
                    }
                )
                index += 1

        suggested = suggest_map(layouts)
        for entry in layouts:
            del entry["counts"]  # internal scoring detail, not part of the contract

        print(
            json.dumps(
                {
                    "ok": True,
                    "slideWidthEmu": int(prs.slide_width),
                    "slideHeightEmu": int(prs.slide_height),
                    "layouts": layouts,
                    "suggestedMap": suggested,
                }
            )
        )
    except Exception as e:  # noqa: BLE001
        fail("introspect_failed", f"introspection failed: {e}")


if __name__ == "__main__":
    main()
